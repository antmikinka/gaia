#!/usr/bin/env python3
"""
GAIA V2 Technical Brief PowerPoint Converter
Converts GAIA_TECHNICAL_BRIEF_v2.md into a comprehensive PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# =============================================================================
# STYLE CONFIGURATION - Bold Signal Theme
# =============================================================================

class StyleConfig:
    """Bold Signal visual style configuration"""
    # Colors
    BG_PRIMARY = RGBColor(26, 26, 26)      # #1a1a1a
    BG_GRADIENT = RGBColor(45, 45, 45)     # #2d2d2d
    CARD_BG = RGBColor(255, 87, 34)        # #FF5722 - Orange accent
    TEXT_PRIMARY = RGBColor(255, 255, 255) # #ffffff
    TEXT_SECONDARY = RGBColor(176, 176, 176) # #b0b0b0
    TEXT_ON_CARD = RGBColor(26, 26, 26)    # #1a1a1a
    ACCENT_LIGHT = RGBColor(255, 138, 101) # #FF8A65

    # Typography
    TITLE_FONT = 'Arial Black'
    BODY_FONT = 'Segoe UI'
    CODE_FONT = 'Consolas'

# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def set_background(slide, prs):
    """Apply dark gradient background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = StyleConfig.BG_PRIMARY

def add_title_slide(prs, title, subtitle, author=""):
    """Create a title slide with Bold Signal styling"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_background(slide, prs)

    # Slide number
    slide_num = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(1), Inches(0.5))
    slide_num.text_frame.text = "01"
    slide_num.text_frame.paragraphs[0].font.size = Pt(72)
    slide_num.text_frame.paragraphs[0].font.color.rgb = StyleConfig.CARD_BG
    slide_num.text_frame.paragraphs[0].font.name = StyleConfig.TITLE_FONT
    slide_num.text_frame.paragraphs[0].font.bold = True
    slide_num.fill.solid()
    slide_num.fill.fore_color.rgb = StyleConfig.BG_PRIMARY
    slide_num.line.fill.background()

    # Title card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(2.5), Inches(8), Inches(2.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = StyleConfig.CARD_BG
    card.line.fill.background()

    # Title text
    tf = card.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.color.rgb = StyleConfig.TEXT_ON_CARD
    p.font.name = StyleConfig.TITLE_FONT
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = StyleConfig.TEXT_SECONDARY
    p.font.name = StyleConfig.BODY_FONT
    p.alignment = PP_ALIGN.CENTER

    if author:
        author_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
        tf = author_box.text_frame
        p = tf.paragraphs[0]
        p.text = author
        p.font.size = Pt(14)
        p.font.color.rgb = StyleConfig.TEXT_SECONDARY
        p.font.name = StyleConfig.BODY_FONT
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_items, slide_number=2):
    """Create a content slide with title and bullet points or structured content"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, prs)

    # Slide number
    slide_num = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(1), Inches(0.5))
    slide_num.text_frame.text = f"{slide_number:02d}"
    slide_num.text_frame.paragraphs[0].font.size = Pt(72)
    slide_num.text_frame.paragraphs[0].font.color.rgb = StyleConfig.CARD_BG
    slide_num.text_frame.paragraphs[0].font.name = StyleConfig.TITLE_FONT
    slide_num.text_frame.paragraphs[0].font.bold = True

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = StyleConfig.TEXT_PRIMARY
    p.font.name = StyleConfig.TITLE_FONT
    p.font.bold = True

    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, dict):
            p.text = item.get('text', '')
            p.font.size = Pt(item.get('size', 18))
            if item.get('bold', False):
                p.font.bold = True
            p.font.color.rgb = StyleConfig.TEXT_PRIMARY
        else:
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = StyleConfig.TEXT_PRIMARY

        p.font.name = StyleConfig.BODY_FONT
        p.space_after = Pt(12)

    return slide

def add_table_slide(prs, title, headers, rows, slide_number=2):
    """Create a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, prs)

    # Slide number
    slide_num = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(1), Inches(0.5))
    slide_num.text_frame.text = f"{slide_number:02d}"
    slide_num.text_frame.paragraphs[0].font.size = Pt(72)
    slide_num.text_frame.paragraphs[0].font.color.rgb = StyleConfig.CARD_BG
    slide_num.text_frame.paragraphs[0].font.name = StyleConfig.TITLE_FONT
    slide_num.text_frame.paragraphs[0].font.bold = True

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = StyleConfig.TEXT_PRIMARY
    p.font.name = StyleConfig.TITLE_FONT
    p.font.bold = True

    # Table
    rows_count = len(rows) + 1  # +1 for header
    cols_count = len(headers)

    table = slide.shapes.add_table(
        rows_count, cols_count,
        Inches(0.5), Inches(1.3),
        Inches(9), Inches(0.8 * rows_count)
    ).table

    # Set column widths
    for i in range(cols_count):
        table.columns[i].width = Inches(9 / cols_count)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = StyleConfig.CARD_BG

        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = StyleConfig.TEXT_ON_CARD
            paragraph.font.name = StyleConfig.BODY_FONT
            paragraph.font.bold = True

    # Data rows
    for row_idx, row in enumerate(rows, 1):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = StyleConfig.TEXT_PRIMARY
                paragraph.font.name = StyleConfig.BODY_FONT

    return slide

def add_diagram_slide(prs, title, diagram_description, slide_number=2):
    """Create a slide with a diagram description (text-based for now)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, prs)

    # Slide number
    slide_num = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(1), Inches(0.5))
    slide_num.text_frame.text = f"{slide_number:02d}"
    slide_num.text_frame.paragraphs[0].font.size = Pt(72)
    slide_num.text_frame.paragraphs[0].font.color.rgb = StyleConfig.CARD_BG
    slide_num.text_frame.paragraphs[0].font.name = StyleConfig.TITLE_FONT
    slide_num.text_frame.paragraphs[0].font.bold = True

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = StyleConfig.TEXT_PRIMARY
    p.font.name = StyleConfig.TITLE_FONT
    p.font.bold = True

    # Diagram description (code-style)
    diagram_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.5))
    tf = diagram_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = diagram_description
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(100, 255, 100)  # Green for code-like appearance
    p.font.name = StyleConfig.CODE_FONT

    return slide

def add_grid_slide(prs, title, cards, slide_number=2):
    """Create a slide with a grid of cards"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, prs)

    # Slide number
    slide_num = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(1), Inches(0.5))
    slide_num.text_frame.text = f"{slide_number:02d}"
    slide_num.text_frame.paragraphs[0].font.size = Pt(72)
    slide_num.text_frame.paragraphs[0].font.color.rgb = StyleConfig.CARD_BG
    slide_num.text_frame.paragraphs[0].font.name = StyleConfig.TITLE_FONT
    slide_num.text_frame.paragraphs[0].font.bold = True

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = StyleConfig.TEXT_PRIMARY
    p.font.name = StyleConfig.TITLE_FONT
    p.font.bold = True

    # Create cards grid (2x2 or 2x3)
    card_width = 3.8
    card_height = 2.2
    start_x = 0.7
    start_y = 1.3

    for idx, card in enumerate(cards):
        row = idx // 2
        col = idx % 2

        x = start_x + col * (card_width + 0.3)
        y = start_y + row * (card_height + 0.3)

        # Card background
        card_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(card_width), Inches(card_height)
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = StyleConfig.CARD_BG
        card_shape.line.fill.background()

        # Card text
        tf = card_shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = card.get('title', '')
        p.font.size = Pt(14)
        p.font.color.rgb = StyleConfig.TEXT_ON_CARD
        p.font.name = StyleConfig.TITLE_FONT
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        if card.get('content'):
            for item in card['content']:
                p_new = tf.add_paragraph()
                p_new.text = f"• {item}"
                p_new.font.size = Pt(11)
                p_new.font.color.rgb = StyleConfig.TEXT_ON_CARD
                p_new.font.name = StyleConfig.BODY_FONT

    return slide

# =============================================================================
# MAIN PRESENTATION CREATION
# =============================================================================

def create_gaia_v2_presentation():
    """Create the complete GAIA V2 Technical Brief PowerPoint"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_num = 1

    # ==========================================================================
    # SLIDE 1: Title Slide
    # ==========================================================================
    add_title_slide(
        prs,
        "GAIA V2",
        "Comprehensive Technical Brief & Implementation Proposal",
        "Anthony Mikinka | March 23, 2026 | github.com/antmikinka"
    )
    slide_num += 1

    # ==========================================================================
    # SLIDE 2: Executive Summary
    # ==========================================================================
    add_content_slide(prs, "Executive Summary", [
        {"text": "GAIA V2 (Generalized Agent Intelligence Architecture)", "size": 20, "bold": True},
        {"text": "Production-proven multi-agent orchestration system"},
        {"text": "Delivers 'one prompt → complete software feature' capability"},
        {"text": "", "size": 8},
        {"text": "Proven Foundation:", "size": 18, "bold": True},
        {"text": "• 8 active production hooks in Safe Haven"},
        {"text": "• 99.8% test pass rate (1120/1122)"},
        {"text": "• 41.8K+ stars (BMAD-METHOD)"},
        {"text": "• 27 quality validation categories"},
        {"text": "• 17 agent specialists across 4 categories"},
        {"text": "• 8 pre-configured pipeline templates"},
    ], slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 3: Key Metrics Table
    # ==========================================================================
    headers = ["Metric", "Status", "Verification"]
    rows = [
        ["Production Hooks", "8 active", "Running in Claude Code"],
        ["Test Pass Rate", "99.8% (1120/1122)", "Verified"],
        ["Community Validation", "41.8K+ stars", "GitHub BMAD-METHOD"],
        ["Quality Categories", "27 validators", "Implemented"],
        ["Agent Specialists", "17 agents", "4 categories"],
        ["Pipeline Templates", "8 configurations", "STANDARD-RAPID-ENTERPRISE-etc"],
        ["Development Timeline", "20 weeks", "Proposed"],
        ["Expected Efficiency Gain", "10x", "With Ryzen AI"],
    ]
    add_table_slide(prs, "Key Metrics", headers, rows, slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 4: Problem Statement
    # ==========================================================================
    headers = ["Pain Point", "Current Solution", "Market Impact"]
    rows = [
        ["Agent Creation Complexity", "Manual coding (weeks)", "$15B (AI dev tools)"],
        ["Quality Assurance Gaps", "Basic or none", "$60B (testing/QA)"],
        ["Cloud Dependency", "100% cloud APIs", "$50B (edge AI)"],
        ["Workflow Fragmentation", "Linear chains", "$25B (orchestration)"],
        ["No Hardware Optimization", "Generic execution", "Missed NPU/GPU potential"],
    ]
    add_table_slide(prs, "Problem Statement - Market Opportunity", headers, rows, slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 5: System Architecture Diagram
    # ==========================================================================
    diagram = """
┌─────────────────────────────────────────────────────────────────────────┐
│                    GAIA V2 SYSTEM ARCHITECTURE                           │
├─────────────────────────────────────────────────────────────────────────┤
│                                                                          │
│  USER GOAL → [PLANNING] → [DEVELOPMENT] → [QUALITY] → [DECISION]       │
│                                      │                                    │
│              ┌───────────────────────┴───────────────────────┐          │
│              │              QUALITY GATE                      │          │
│              │              Score >= Threshold?               │          │
│              │              YES → SHIP ✓                      │          │
│              │              NO  → EXTRACT DEFECTS             │          │
│              │                   LOOP TO PLANNING             │          │
│              │                   (unlimited iterations)       │          │
│              └────────────────────────────────────────────────┘          │
│                                                                          │
│  AGENT CATEGORIES:        PIPELINE TEMPLATES:                            │
│  - PLANNING (4 agents)    - STANDARD (90/100)                            │
│  - DEVELOPMENT (5 agents) - RAPID (75/100)                               │
│  - REVIEW (5 agents)      - ENTERPRISE (95/100)                          │
│  - MANAGEMENT (3 agents)  - 8 total configurations                       │
│                                                                          │
└─────────────────────────────────────────────────────────────────────────┘
"""
    add_diagram_slide(prs, "System Architecture", diagram, slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 6: Core Innovation - Recursive Iterative Pipeline
    # ==========================================================================
    add_content_slide(prs, "Core Innovation: Recursive Iterative Pipeline", [
        {"text": "Key Differentiator:", "size": 20, "bold": True},
        {"text": "No artificial max iterations - continues until quality threshold met"},
        {"text": "", "size": 8},
        {"text": "Pipeline Flow:", "size": 18, "bold": True},
        {"text": "1. PLANNING PHASE - Requirement analysis, Architecture design, Task decomposition, Agent selection"},
        {"text": "2. DEVELOPMENT PHASE - Code generation, Component creation, Integration"},
        {"text": "3. QUALITY PHASE - 27 validation categories, Weighted scoring, Defect identification"},
        {"text": "4. DECISION PHASE - Score >= Threshold? → SHIP ✓ | Score < Threshold? → LOOP"},
        {"text": "", "size": 8},
        {"text": "State-Based Routing:", "size": 18, "bold": True},
        {"text": "• security defect → security-auditor → loop back"},
        {"text": "• performance issue → performance-analyst → fix"},
        {"text": "• API task → api-designer → plan"},
    ], slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 7: Agent Ecosystem Grid
    # ==========================================================================
    cards = [
        {
            "title": "PLANNING (4 agents)",
            "content": [
                "planning-analysis-strategist",
                "solutions-architect",
                "api-designer",
                "database-architect"
            ]
        },
        {
            "title": "DEVELOPMENT (5 agents)",
            "content": [
                "senior-developer",
                "frontend-specialist",
                "backend-specialist",
                "devops-engineer",
                "data-engineer"
            ]
        },
        {
            "title": "REVIEW (5 agents)",
            "content": [
                "quality-reviewer",
                "security-auditor",
                "performance-analyst",
                "accessibility-reviewer",
                "test-coverage-analyzer"
            ]
        },
        {
            "title": "MANAGEMENT (3 agents)",
            "content": [
                "software-program-manager",
                "technical-writer",
                "release-manager"
            ]
        }
    ]
    add_grid_slide(prs, "Agent Ecosystem - 17 Specialists", cards, slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 8: Pipeline Templates Table
    # ==========================================================================
    headers = ["Template", "Threshold", "Agent Sequence", "Use Case"]
    rows = [
        ["STANDARD", "90/100", "Planning → Dev → QA → Manager", "Features, APIs"],
        ["RAPID", "75/100", "Planning → Dev → QA", "Prototypes, MVPs"],
        ["ENTERPRISE", "95/100", "Planning → Dev → QA → Security → Perf → Manager", "Production, Security"],
        ["DOCUMENTATION", "85/100", "Tech Writer → Reviewer → Editor", "API docs, guides"],
        ["TESTING", "90/100", "Test Architect → Dev → QA → Coverage", "Test creation"],
        ["FRONTEND", "88/100", "API Designer → Frontend → QA → Accessibility", "UI components"],
        ["BACKEND", "90/100", "API Designer → Backend → QA → Security", "REST APIs"],
        ["DATA-ML", "88/100", "DB Architect → Data Engineer → QA", "Data pipelines"],
    ]
    add_table_slide(prs, "Pipeline Templates - 8 Configurations", headers, rows, slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 9: Quality Scoring System
    # ==========================================================================
    headers = ["Dimension", "Weight", "Validation Categories"]
    rows = [
        ["Code Quality", "25%", "Syntax, Style, Complexity, DRY, SOLID, Error Handling"],
        ["Requirements Coverage", "25%", "Feature Completeness, Edge Cases, User Stories"],
        ["Testing", "20%", "Unit Tests, Integration Tests, Coverage, Mock Quality"],
        ["Documentation", "15%", "Docstrings, README, API Docs, Comments"],
        ["Best Practices", "15%", "Security, Performance, Accessibility, Maintainability"],
    ]
    add_table_slide(prs, "Quality Scoring System - 27 Categories", headers, rows, slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 10: State-Based Routing Examples
    # ==========================================================================
    add_content_slide(prs, "State-Based Routing - Dynamic Agent Selection", [
        {"text": "Routing Rules (YAML Configuration):", "size": 20, "bold": True},
        {"text": "", "size": 8},
        {"text": "• defect_type == 'security' → security-auditor (mandatory_fix)", "size": 16},
        {"text": "• defect_type == 'performance' AND severity >= 8 → performance-analyst", "size": 16},
        {"text": "• task_type contains 'api' → api-designer", "size": 16},
        {"text": "• task_type contains 'database' → database-architect", "size": 16},
        {"text": "• test_coverage < 90 → test-coverage-analyzer (loop_back)", "size": 16},
        {"text": "• ui_component == true → accessibility-reviewer", "size": 16},
        {"text": "", "size": 8},
        {"text": "Key Benefit:", "size": 18, "bold": True},
        {"text": "Routes to the RIGHT specialist based on defect type, not fixed sequence"},
    ], slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 11: Hook System Architecture
    # ==========================================================================
    add_content_slide(prs, "Hook System - 16 Events, 8 Production Hooks", [
        {"text": "Hook Events:", "size": 20, "bold": True},
        {"text": "• Pipeline Lifecycle (5): init, start, complete, fail, cancel"},
        {"text": "• Phase Events (3): phase_start, phase_complete, phase_fail"},
        {"text": "• Loop Events (3): loop_start, loop_complete, loop_defects"},
        {"text": "• Quality Events (3): quality_eval, threshold_met, threshold_failed"},
        {"text": "• Agent Events (2): agent_invoke, agent_complete"},
        {"text": "", "size": 8},
        {"text": "Production Hooks (8):", "size": 18, "bold": True},
        {"text": "1. PreActionValidationHook  5. QualityGateHook"},
        {"text": "2. PostActionValidationHook 6. DefectExtractionHook"},
        {"text": "3. ContextInjectionHook     7. PipelineNotificationHook"},
        {"text": "4. OutputProcessingHook     8. ChronicleHarvestHook"},
    ], slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 12: Safe Haven Integration
    # ==========================================================================
    diagram = """
┌─────────────────────────────────────────────────────────────────┐
│ HOOK INTEGRATION ARCHITECTURE                                    │
├─────────────────────────────────────────────────────────────────┤
│                                                                  │
│  CLAUDE CODE HOOK SYSTEM (Safe Haven)                            │
│  ├── pre-compaction-validation.py                                │
│  ├── context-preservation-optimizer.py                           │
│  └── post-compaction-monitor.py                                  │
│                                                                  │
│  GAIA HOOK SYSTEM (Extension)                                    │
│  ├── PreActionValidationHook                                     │
│  ├── QualityGateHook                                             │
│  └── DefectExtractionHook                                        │
│                                                                  │
│  INTEGRATION: GAIA hooks execute within Safe Haven context       │
│  - Shared logging infrastructure                                 │
│  - Shared metrics/monitoring                                     │
│  - Shared error handling                                         │
│                                                                  │
└─────────────────────────────────────────────────────────────────┘
"""
    add_diagram_slide(prs, "Safe Haven Hook Integration", diagram, slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 13: Implementation Plan - 20 Weeks
    # ==========================================================================
    headers = ["Phase", "Duration", "Deliverables", "Success Criteria"]
    rows = [
        ["Phase 1", "Weeks 1-4", "Core Pipeline Engine", "State machine, Loop manager,\nDecision engine, Quality scorer"],
        ["Phase 2", "Weeks 5-12", "Productization", "Agent registry, Templates,\nTest suite, Documentation"],
        ["Phase 3", "Weeks 13-20", "AMD Ryzen AI Integration", "NPU optimization,\nChromaDB, MCP server,\nEnterprise pilot"],
    ]
    add_table_slide(prs, "Implementation Plan - 20 Weeks", headers, rows, slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 14: Resource Requirements
    # ==========================================================================
    headers = ["Phase", "Team Size", "Key Roles", "Duration"]
    rows = [
        ["Phase 1", "3-4 FTE", "1 Architect, 2-3 Engineers", "4 weeks"],
        ["Phase 2", "5-6 FTE", "+ QA Engineer, Tech Writer", "8 weeks"],
        ["Phase 3", "4-5 FTE", "+ AMD Liaison", "8 weeks"],
        ["TOTAL", "10 FTE", "Cross-functional team", "20 weeks"],
    ]
    add_table_slide(prs, "Resource Requirements", headers, rows, slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 15: Risk Analysis
    # ==========================================================================
    headers = ["Risk", "Probability", "Impact", "Mitigation"]
    rows = [
        ["Quality gates fail at scale", "Low", "High", "Continuous validation, enterprise pilots"],
        ["Ryzen AI underperforms", "Low", "Medium", "Multi-hardware support roadmap"],
        ["Hook system doesn't generalize", "Medium", "High", "Document patterns early, test across domains"],
        ["AI winter / hype crash", "Low", "High", "Focus on enterprise ROI, not hype"],
        ["Developer resistance to AI", "Low", "Medium", "Bottom-up adoption (dev-first)"],
    ]
    add_table_slide(prs, "Risk Analysis & Mitigation", headers, rows, slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 16: Business Case - ROI Analysis
    # ==========================================================================
    headers = ["Metric", "Current", "With GAIA", "Improvement"]
    rows = [
        ["Dev cost per feature", "$10K", "$1K", "90% reduction"],
        ["Time to market", "6 months", "2 weeks", "12x faster"],
        ["Startup capital needed", "$2M", "$200K", "90% reduction"],
        ["Enterprise dev team ROI", "1x", "10x", "10x improvement"],
    ]
    slide = add_table_slide(prs, "Business Case - ROI Analysis", headers, rows, slide_num)

    # Add economic value callout
    econ_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(5.5), Inches(8), Inches(1.2)
    )
    econ_box.fill.solid()
    econ_box.fill.fore_color.rgb = StyleConfig.CARD_BG
    econ_box.line.fill.background()

    tf = econ_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Economic Value: $650K/year/team × 1M teams = $650B annual productivity gain"
    p.font.size = Pt(18)
    p.font.color.rgb = StyleConfig.TEXT_ON_CARD
    p.font.name = StyleConfig.TITLE_FONT
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    slide_num += 1

    # ==========================================================================
    # SLIDE 17: Revenue Projection
    # ==========================================================================
    headers = ["Revenue Stream", "Year 1", "Year 3", "Year 5"]
    rows = [
        ["Enterprise Licenses", "$2.5M", "$40M", "$120M"],
        ["Template Marketplace (15%)", "$0.5M", "$8M", "$25M"],
        ["AMD Hardware Revenue Share", "$0", "$12M", "$50M"],
        ["Training/Certification", "$0", "$5M", "$15M"],
        ["TOTAL", "$3M", "$65M", "$210M"],
    ]
    add_table_slide(prs, "Revenue Projection - 5 Year Outlook", headers, rows, slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 18: AMD Partnership
    # ==========================================================================
    add_content_slide(prs, "AMD Partnership - 'Intel Inside' Moment", [
        {"text": "What AMD Gets:", "size": 20, "bold": True},
        {"text": "• Differentiated Ryzen AI value proposition ('GAIA-Optimized' badge)"},
        {"text": "• Developer mindshare: 500K+ by Year 3"},
        {"text": "• Hardware sales driver: $500M+ by Year 5"},
        {"text": "• Strategic positioning vs Intel/NVIDIA"},
        {"text": "", "size": 8},
        {"text": "Partnership Requirements:", "size": 18, "bold": True},
        {"text": "• Hardware Access: Ryzen AI dev kits → NPU optimization"},
        {"text": "• Technical Liaison: AMD engineer (part-time) → Architecture guidance"},
        {"text": "• Marketing Support: Co-marketing budget → Developer outreach"},
        {"text": "• Enterprise Introductions: Customer pilot contacts → First deployments"},
    ], slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 19: Next Steps
    # ==========================================================================
    add_content_slide(prs, "Next Steps - Immediate Actions (Week 1)", [
        {"text": "1. Team Allocation", "size": 18, "bold": True},
        {"text": "   Finalize 3-4 FTE for Phase 1"},
        {"text": "", "size": 6},
        {"text": "2. AMD Partnership Agreement", "size": 18, "bold": True},
        {"text": "   Sign technical partnership"},
        {"text": "", "size": 6},
        {"text": "3. Development Infrastructure", "size": 18, "bold": True},
        {"text": "   Set up repositories, CI/CD"},
        {"text": "", "size": 6},
        {"text": "4. Phase 1 Kickoff", "size": 18, "bold": True},
        {"text": "   Begin Week 1 implementation"},
        {"text": "", "size": 8},
        {"text": "Go/No-Go Gates: Week 4 | Week 12 | Week 20", "size": 16, "bold": True},
    ], slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 20: Implementation Status
    # ==========================================================================
    headers = ["Phase", "Component", "Status", "Tests"]
    rows = [
        ["Phase 1", "Pipeline State Machine", "✅ Complete", "24 passing"],
        ["Phase 1", "Loop Manager", "✅ Complete", "19 passing"],
        ["Phase 1", "Decision Engine", "✅ Complete", "17 passing"],
        ["Phase 1", "Quality Scorer", "✅ Complete", "23 passing"],
        ["Phase 1", "Agent Registry", "✅ Complete", "12 passing"],
        ["Phase 1", "Hook System", "✅ Complete", "8 passing"],
        ["TOTAL", "", "", "103 passing"],
    ]
    add_table_slide(prs, "Implementation Status - Phase 1 Complete", headers, rows, slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 21: File Structure
    # ==========================================================================
    file_structure = """
gaia-proposal/
├── GAIA_TECHNICAL_BRIEF_v2.md      # Comprehensive V2 specification
├── GAIA_VISION_2030.md              # Long-term vision
├── GAIA_IMPLEMENTATION_PLAN.pptx    # Implementation PowerPoint
├── GAIA_IMPLEMENTATION_PLAN.html    # HTML presentation
├── README.md                        # Project overview
├── gaia/                            # GAIA Core Pipeline Engine
│   ├── src/gaia/
│   │   ├── pipeline/               # State machine, loops, decisions
│   │   ├── quality/                # Quality scorer, 27 validators
│   │   ├── agents/                 # Agent registry, 17 definitions
│   │   ├── hooks/                  # Hook system, 8 production hooks
│   │   └── utils/                  # Utilities, logging, exceptions
│   ├── config/agents/              # Agent YAML configurations
│   ├── tests/                      # 103 passing tests
│   └── pyproject.toml
└── images/                          # Assets
"""
    add_diagram_slide(prs, "Repository File Structure", file_structure, slide_num)
    slide_num += 1

    # ==========================================================================
    # SLIDE 22: Contact & Approval
    # ==========================================================================
    add_content_slide(prs, "Contact & Next Steps", [
        {"text": "Anthony Mikinka", "size": 24, "bold": True},
        {"text": "Email: anthony.mikinka@gmail.com", "size": 18},
        {"text": "GitHub: github.com/antmikinka", "size": 18},
        {"text": "Location: Budapest, Hungary", "size": 18},
        {"text": "", "size": 8},
        {"text": "Repository:", "size": 18, "bold": True},
        {"text": "https://github.com/antmikinka/gaia-proposal.git", "size": 16},
        {"text": "", "size": 8},
        {"text": "Document Version: 2.0", "size": 14},
        {"text": "Last Updated: March 23, 2026", "size": 14},
    ], slide_num)

    return prs

if __name__ == "__main__":
    print("Creating GAIA V2 Technical Brief PowerPoint...")
    prs = create_gaia_v2_presentation()

    output_file = "GAIA_TECHNICAL_BRIEF_v2.pptx"
    prs.save(output_file)
    print(f"Presentation saved to: {output_file}")
    print(f"Total slides: {len(prs.slides)}")
