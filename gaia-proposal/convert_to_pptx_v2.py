#!/usr/bin/env python3
"""Convert GAIA HTML presentation to PowerPoint format - 21 slides exactly"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Brand colors
BRAND_ORANGE = RGBColor(255, 87, 34)
BRAND_DARK = RGBColor(26, 26, 26)
BRAND_WHITE = RGBColor(255, 255, 255)
BRAND_GRAY = RGBColor(45, 45, 45)

def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_DARK

def add_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12), Inches(2))
    p = tb.text_frame.paragraphs[0]
    p.text = "GAIA"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(72)
    p.font.color.rgb = BRAND_WHITE

    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(1.5))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Generalized Agent Intelligence Architecture"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(28)
    p2.font.color.rgb = BRAND_ORANGE

    # Tagline
    tb3 = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12), Inches(1))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = '"One Prompt -> Complete Software Feature"'
    p3.font.name = 'Space Grotesk'
    p3.font.size = Pt(22)
    p3.font.color.rgb = BRAND_WHITE

    # Footer
    tb4 = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12), Inches(1))
    p4 = tb4.text_frame.paragraphs[0]
    p4.text = "Anthony Mikinka | github.com/antmikinka\nProductizing Battle-Tested Hook Architecture for AMD Ryzen AI"
    p4.font.name = 'Space Grotesk'
    p4.font.size = Pt(16)
    p4.font.color.rgb = BRAND_WHITE

    return slide

def add_executive_summary():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "Executive Summary"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(40)
    p.font.color.rgb = BRAND_ORANGE

    # Stats (3 columns)
    stat_w, stat_h = Inches(3.5), Inches(1.8)
    for i, (num, lbl) in enumerate([
        ('41.8K+', 'BMAD-METHOD Stars\n(Agent Templates Contributor)'),
        ('99.8%', 'Test Pass Rate\n(1120/1122 tests)'),
        ('20 Weeks', 'To Production-Ready\n(Phased Delivery)')
    ]):
        x = Inches(0.5) + i * (stat_w + Inches(0.3))
        # Card bg
        card = slide.shapes.add_shape(1, x, Inches(1.2), stat_w, stat_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 100, 50)
        card.line.fill.background()
        # Number
        tb_num = slide.shapes.add_textbox(x, Inches(1.3), stat_w, Inches(0.9))
        p_num = tb_num.text_frame.paragraphs[0]
        p_num.text = num
        p_num.font.name = 'Archivo Black'
        p_num.font.size = Pt(48)
        p_num.font.color.rgb = BRAND_WHITE
        p_num.alignment = PP_ALIGN.CENTER
        # Label
        tb_lbl = slide.shapes.add_textbox(x, Inches(2.2), stat_w, Inches(0.7))
        p_lbl = tb_lbl.text_frame.paragraphs[0]
        p_lbl.text = lbl
        p_lbl.font.name = 'Space Grotesk'
        p_lbl.font.size = Pt(13)
        p_lbl.font.color.rgb = BRAND_WHITE
        p_lbl.alignment = PP_ALIGN.CENTER

    # Opportunity statement
    tb_opp = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12), Inches(0.8))
    p_opp = tb_opp.text_frame.paragraphs[0]
    p_opp.text = "The Opportunity: Productize proven 'one prompt -> feature' capability for AMD Ryzen AI"
    p_opp.font.name = 'Space Grotesk'
    p_opp.font.size = Pt(18)
    p_opp.font.color.rgb = BRAND_WHITE
    p_opp.font.bold = True

    # Table
    table = slide.shapes.add_table(4, 3, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.8)).table
    headers = ['Metric', 'Current', 'With GAIA + Ryzen AI']
    rows = [
        ['Development Speed', 'Hours-long automated', '10x faster (local NPU)'],
        ['Quality Threshold', '90-95/100', 'Same + hardware optimization'],
        ['Cloud Dependency', '100% cloud LLMs', 'Local execution']
    ]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_ORANGE
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Archivo Black'
            p.font.size = Pt(12)
            p.font.color.rgb = BRAND_WHITE
            p.font.bold = True
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.name = 'Space Grotesk'
                p.font.size = Pt(13)
                p.font.color.rgb = BRAND_WHITE

    return slide

def add_safe_haven():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = "Proof: Safe Haven Running in Claude Code"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Production Hook System - Live Integration"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    # Code block
    code_bg = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.5))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = BRAND_GRAY
    code_bg.line.fill.background()

    code = """SAFE HAVEN: 8 PRODUCTION HOOKS
    pre-compaction-validation.py (28KB) - Critical context analysis
    context-preservation-optimizer.py (117KB) - 5 compression strategies
    post-compaction-monitor.py (50KB) - Real-time metrics
    quality-assurance-validator.py - 25+ validation categories
    compliance-check.py - Regulatory compliance
    content-quality-monitor.py - Content validation
    performance-validation.py - Performance gates

    EXECUTION: Automatically triggered during Claude Code sessions
    STATUS: ACTIVE and running NOW"""

    tb_code = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(4.1))
    tf = tb_code.text_frame
    tf.word_wrap = True
    p_code = tf.paragraphs[0]
    p_code.text = code
    p_code.font.name = 'Courier New'
    p_code.font.size = Pt(14)
    p_code.font.color.rgb = BRAND_WHITE

    # Footer note
    tb_note = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12), Inches(0.5))
    p_note = tb_note.text_frame.paragraphs[0]
    p_note.text = "This is MY architecture. GAIA productizes it."
    p_note.font.name = 'Space Grotesk'
    p_note.font.size = Pt(16)
    p_note.font.color.rgb = BRAND_ORANGE
    p_note.font.bold = True

    return slide

def add_recursive_pipeline():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = "The Core: Recursive Iterative Pipeline"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Quality-Gated Loop Architecture"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    code_bg = slide.shapes.add_shape(1, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.3))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = BRAND_GRAY

    code = """+---------------------------------------------------------------------+
|         RECURSIVE ITERATIVE PIPELINE (Auto-Pilot)                   |
+---------------------------------------------------------------------+
|  USER GOAL -> PLANNING -> DEVELOPMENT -> QUALITY -> DECISION        |
|                                 |                                   |
|         +-----------------------+----------------------------+      |
|         |                  QUALITY GATE                       |      |
|         |   Score >= 90?                                       |      |
|         |   YES -> @software-program-manager -> SHIP          |      |
|         |   NO  -> Defects -> Loop back to PLANNING           |      |
|         |        (Unlimited iterations until quality met)      |      |
|         +------------------------------------------------------+      |
+---------------------------------------------------------------------+

    KEY INNOVATION: No artificial max iterations. Continues until quality >= threshold.
    Source: RECURSIVE-ITERATIVE-PIPELINE.md"""

    tb_code = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.1))
    tf = tb_code.text_frame
    tf.word_wrap = True
    p_code = tf.paragraphs[0]
    p_code.text = code
    p_code.font.name = 'Courier New'
    p_code.font.size = Pt(10)
    p_code.font.color.rgb = BRAND_WHITE

    return slide

def add_pipeline_templates():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = "Pipeline Templates (Already Built)"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Configurable Quality Thresholds"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    # Table
    table = slide.shapes.add_table(5, 4, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.0)).table
    headers = ['Template', 'Quality Threshold', 'Agent Sequence', 'Use Case']
    rows = [
        ['STANDARD', '90/100', 'Planning -> Dev -> QA -> Manager', 'Features, APIs'],
        ['RAPID', '75/100', 'Planning -> Dev -> QA', 'Prototypes, MVPs'],
        ['ENTERPRISE', '95/100', 'Planning -> Dev -> QA -> Security -> Perf -> Mgr', 'Production, Security'],
        ['DOCUMENTATION', '85/100', 'Tech Writer -> Reviewer -> Editor', 'API docs, guides']
    ]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_ORANGE
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Archivo Black'
            p.font.size = Pt(10)
            p.font.color.rgb = BRAND_WHITE
            p.font.bold = True
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.name = 'Space Grotesk'
                p.font.size = Pt(10)
                p.font.color.rgb = BRAND_WHITE

    # Quality weights
    weights = """QUALITY WEIGHTS: Code Quality 25% | Requirements 25% | Testing 20% | Documentation 15% | Best Practices 15%"""
    tb_w = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.6))
    p_w = tb_w.text_frame.paragraphs[0]
    p_w.text = weights
    p_w.font.name = 'Space Grotesk'
    p_w.font.size = Pt(14)
    p_w.font.color.rgb = BRAND_WHITE

    tb_src = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12), Inches(0.4))
    p_src = tb_src.text_frame.paragraphs[0]
    p_src.text = "Source: auto-pilot-templates.yml, recursive-pipeline-templates.yml"
    p_src.font.name = 'Space Grotesk'
    p_src.font.size = Pt(11)
    p_src.font.color.rgb = RGBColor(150, 150, 150)

    return slide

def add_agent_categories():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "Agent Categories (State-Based Routing)"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(30)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Multi-Agent Selection System"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(16)
    p2.font.color.rgb = BRAND_ORANGE

    code = """AGENT CATEGORIES:
    PLANNING:       planning-analysis-strategist | solutions-architect | api-designer | database-architect
    DEVELOPMENT:    senior-developer | frontend-specialist | backend-specialist | devops-engineer | data-engineer
    REVIEW:         quality-reviewer | security-auditor | performance-analyst | accessibility-reviewer | test-coverage-analyzer
    MANAGEMENT:     software-program-manager | technical-writer | release-manager

    AUTO-SELECTION: Based on task triggers ("api" -> api-designer, "security" -> security-auditor)
    Source: recursive-pipeline-templates.yml"""

    code_bg = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.8))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = BRAND_GRAY

    tb_code = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.6))
    tf = tb_code.text_frame
    tf.word_wrap = True
    p_code = tf.paragraphs[0]
    p_code.text = code
    p_code.font.name = 'Courier New'
    p_code.font.size = Pt(11)
    p_code.font.color.rgb = BRAND_WHITE

    return slide

def add_3d_cube():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "3D Cube Matrix Classification"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Multi-Dimensional Agent Organization"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    code = """                PHASES (Z-axis): Research -> Planning -> Development -> QA -> Feedback
                          /
                         /   COMPONENT TYPE (Y-axis): tasks | checklists | tools | workflows | data
                        /
                AGENT PERSONA (X-axis): planning | developer | qa | researcher | manager

    EXAMPLE: execute-checklist.md -> X:qa-reviewer, Y:checklist, Z:QA phase

    Enables: Intelligent agent routing | Component type awareness | Cross-phase state injection
    Source: ORCHESTRATION-SPEC-v1.md (92KB specification)"""

    code_bg = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.8))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = BRAND_GRAY

    tb_code = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.6))
    tf = tb_code.text_frame
    tf.word_wrap = True
    p_code = tf.paragraphs[0]
    p_code.text = code
    p_code.font.name = 'Courier New'
    p_code.font.size = Pt(12)
    p_code.font.color.rgb = BRAND_WHITE

    return slide

def add_orchestration_modes():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "Orchestration Modes"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Configurable Automation Levels"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    # Table
    table = slide.shapes.add_table(4, 4, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.2)).table
    headers = ['Mode', 'Automation', 'User Control', 'Use Case']
    rows = [
        ['Manual', 'User invokes each agent', 'Full control', 'Exploratory, debugging'],
        ['Guided', 'System suggests, user confirms', 'Approval required', 'Standard workflows (DEFAULT)'],
        ['Autonomous', 'Auto-invokes with gates', 'Phase boundaries only', 'Repeatable pipelines']
    ]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_ORANGE
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Archivo Black'
            p.font.size = Pt(11)
            p.font.color.rgb = BRAND_WHITE
            p.font.bold = True
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.name = 'Space Grotesk'
                p.font.size = Pt(11)
                p.font.color.rgb = BRAND_WHITE

    # Execution patterns
    patterns = """EXECUTION PATTERNS: Linear Pipeline | Iterative Spiral | Parallel Execution | Conditional Branching"""
    tb_p = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.6))
    p_p = tb_p.text_frame.paragraphs[0]
    p_p.text = patterns
    p_p.font.name = 'Space Grotesk'
    p_p.font.size = Pt(14)
    p_p.font.color.rgb = BRAND_WHITE

    return slide

def add_nexus_architecture():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "Nexus: 4-Layer Architecture"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Complete Orchestration Stack"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    code = """+--------------------------------------------------+
|  USER / CLAUDE CODE (Natural Conversation)       |
|  @agent mention -> automatic routing             |
+--------------------------------------------------+
|  ORCHESTRATION LAYER                             |
|  OrchestratorAgent | PhaseController | Feedback  |
|  ConfigLoader      | StateManager    | Metrics   |
+--------------------------------------------------+
|  RUNTIME INFRASTRUCTURE (Nexus)                  |
|  PipelinePlanner | ContextInjector | AgentRegistry|
|  Chronicle (temporal) | Workspace (spatial)      |
+--------------------------------------------------+
|  BUILD-TIME INFRASTRUCTURE                       |
|  SessionManager | Agent Templates | Hook System  |
+--------------------------------------------------+
|  CLAUDE CODE HOST                                |
|  Task Tool | File Ops | MCP Servers | ChromaDB   |
+--------------------------------------------------+
    Source: architecture-diagram.html"""

    code_bg = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.0))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = BRAND_GRAY

    tb_code = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.8))
    tf = tb_code.text_frame
    tf.word_wrap = True
    p_code = tf.paragraphs[0]
    p_code.text = code
    p_code.font.name = 'Courier New'
    p_code.font.size = Pt(10)
    p_code.font.color.rgb = BRAND_WHITE

    return slide

def add_bmad():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "Track Record: BMAD-METHOD"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(36)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Agent Templates Before Viral Growth"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    # Stats
    stat_w, stat_h = Inches(3.5), Inches(1.8)
    for i, (num, lbl) in enumerate([
        ('<10K', 'Stars When I Contributed\n(Early Stage)'),
        ('41.8K+', 'Stars Now\n(Viral Growth)'),
        ('Still Ref.', 'My Templates\n(Lasting Impact)')
    ]):
        x = Inches(0.5) + i * (stat_w + Inches(0.3))
        card = slide.shapes.add_shape(1, x, Inches(1.3), stat_w, stat_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 100, 50)
        card.line.fill.background()

        tb_num = slide.shapes.add_textbox(x, Inches(1.4), stat_w, Inches(0.9))
        p_num = tb_num.text_frame.paragraphs[0]
        p_num.text = num
        p_num.font.name = 'Archivo Black'
        p_num.font.size = Pt(44)
        p_num.font.color.rgb = BRAND_WHITE
        p_num.alignment = PP_ALIGN.CENTER

        tb_lbl = slide.shapes.add_textbox(x, Inches(2.3), stat_w, Inches(0.7))
        p_lbl = tb_lbl.text_frame.paragraphs[0]
        p_lbl.text = lbl
        p_lbl.font.name = 'Space Grotesk'
        p_lbl.font.size = Pt(13)
        p_lbl.font.color.rgb = BRAND_WHITE
        p_lbl.alignment = PP_ALIGN.CENTER

    # Statement
    tb_stmt = slide.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(12), Inches(1.0))
    p_stmt = tb_stmt.text_frame.paragraphs[0]
    p_stmt.text = "BMAD-METHOD (github.com/bmad-code-org/BMAD-METHOD)\nCreated agent templates before the project hit mainstream. Still referenced today."
    p_stmt.font.name = 'Space Grotesk'
    p_stmt.font.size = Pt(16)
    p_stmt.font.color.rgb = BRAND_WHITE

    tb_final = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12), Inches(0.6))
    p_final = tb_final.text_frame.paragraphs[0]
    p_final.text = "I don't just design systems. I BUILD them before they're cool."
    p_final.font.name = 'Space Grotesk'
    p_final.font.size = Pt(18)
    p_final.font.color.rgb = BRAND_ORANGE
    p_final.font.bold = True

    return slide

def add_metrics():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "Metrics & Performance"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(36)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Real Numbers from Production"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    # Stats
    stat_w, stat_h = Inches(3.5), Inches(1.6)
    for i, (num, lbl) in enumerate([
        ('33/33', 'Auto-Pilot Tests\n100% Passing'),
        ('1120/1122', 'Full Test Suite\n99.8% Passing'),
        ('8', 'Core Hook Files\nProduction Ready')
    ]):
        x = Inches(0.5) + i * (stat_w + Inches(0.3))
        card = slide.shapes.add_shape(1, x, Inches(1.2), stat_w, stat_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 100, 50)
        card.line.fill.background()

        tb_num = slide.shapes.add_textbox(x, Inches(1.3), stat_w, Inches(0.8))
        p_num = tb_num.text_frame.paragraphs[0]
        p_num.text = num
        p_num.font.name = 'Archivo Black'
        p_num.font.size = Pt(44)
        p_num.font.color.rgb = BRAND_WHITE
        p_num.alignment = PP_ALIGN.CENTER

        tb_lbl = slide.shapes.add_textbox(x, Inches(2.1), stat_w, Inches(0.6))
        p_lbl = tb_lbl.text_frame.paragraphs[0]
        p_lbl.text = lbl
        p_lbl.font.name = 'Space Grotesk'
        p_lbl.font.size = Pt(12)
        p_lbl.font.color.rgb = BRAND_WHITE
        p_lbl.alignment = PP_ALIGN.CENTER

    # Additional info
    doc_info = """DOCUMENTATION: 8 categorized docs | 4 visual diagrams (2 SVG, 1 HTML, 1 ASCII) | Implementation guide | 5+ pipeline templates"""
    tb_doc = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.8))
    p_doc = tb_doc.text_frame.paragraphs[0]
    p_doc.text = doc_info
    p_doc.font.name = 'Space Grotesk'
    p_doc.font.size = Pt(15)
    p_doc.font.color.rgb = BRAND_WHITE

    return slide

def add_use_cases():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "Real-World Use Cases"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "What This Can Build"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    code = """EXAMPLE WORKFLOWS (Already Proven):

    1. PATENT ECOSYSTEM (CCMG):
       - 12+ specialist agents | 780-line orchestrator | 25+ legal templates | 8 commands

    2. RESEARCH AGENT 2:
       - 65+ components | 4 Python hooks | 5 QA checklists | MCP server configs

    3. NEXUS AGENT ORCHESTRATOR:
       - 12+ specialist agents | Multi-layer testing | MCP integration | Self-configuring

    All generated from prompts. All production-ready."""

    code_bg = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.8))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = BRAND_GRAY

    tb_code = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.6))
    tf = tb_code.text_frame
    tf.word_wrap = True
    p_code = tf.paragraphs[0]
    p_code.text = code
    p_code.font.name = 'Courier New'
    p_code.font.size = Pt(12)
    p_code.font.color.rgb = BRAND_WHITE

    return slide

def add_market_opportunity():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "Market Opportunity"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Enterprise AI Development Pain Points"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    # Table
    table = slide.shapes.add_table(5, 4, Inches(0.5), Inches(1.2), Inches(12.3), Inches(3.2)).table
    headers = ['Pain Point', 'Current', 'GAIA Solution', 'Market Size']
    rows = [
        ['Agent Creation', 'Weeks manual coding', 'One prompt -> ecosystem', '$15B (AI dev tools)'],
        ['Quality Assurance', 'Basic or none', '25+ validation categories', '$60B (testing/QA)'],
        ['Cloud Dependency', '100% cloud APIs', 'Local Ryzen AI execution', '$50B (edge AI)'],
        ['Workflow Complexity', 'Linear chains', 'Recursive quality-gated', '$25B (orchestration)']
    ]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_ORANGE
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Archivo Black'
            p.font.size = Pt(10)
            p.font.color.rgb = BRAND_WHITE
            p.font.bold = True
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.name = 'Space Grotesk'
                p.font.size = Pt(10)
                p.font.color.rgb = BRAND_WHITE

    # TAM
    tb_tam = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12), Inches(0.8))
    p_tam = tb_tam.text_frame.paragraphs[0]
    p_tam.text = "Total Addressable Market: $150B+ across AI development, QA, edge computing, and orchestration"
    p_tam.font.name = 'Space Grotesk'
    p_tam.font.size = Pt(18)
    p_tam.font.color.rgb = BRAND_WHITE
    p_tam.font.bold = True

    return slide

def add_competitive():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "Competitive Landscape"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Why GAIA Wins"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    # Table
    table = slide.shapes.add_table(5, 5, Inches(0.5), Inches(1.2), Inches(12.3), Inches(3.0)).table
    headers = ['Platform', 'Agent Creation', 'Quality Gates', 'Hardware Opt.', 'Battle-Tested']
    rows = [
        ['AutoGen', 'Manual', 'Basic', 'None', 'Academic demos'],
        ['CrewAI', 'Manual', 'Basic', 'None', 'Startup pilots'],
        ['LangChain', 'Low-code', 'Limited', 'None', 'Varied'],
        ['GAIA', 'One prompt', '25+ categories', 'Ryzen AI', 'Production']
    ]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_ORANGE
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Archivo Black'
            p.font.size = Pt(10)
            p.font.color.rgb = BRAND_WHITE
            p.font.bold = True
    for ri, row in enumerate(rows, 1):
        highlight = (ri == 4)
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            if highlight:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 100, 50)
            for p in cell.text_frame.paragraphs:
                p.font.name = 'Space Grotesk'
                p.font.size = Pt(11)
                p.font.color.rgb = BRAND_DARK if highlight else BRAND_WHITE
                if highlight:
                    p.font.bold = True

    tb_win = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(12), Inches(0.6))
    p_win = tb_win.text_frame.paragraphs[0]
    p_win.text = "You can copy code. You can't copy PROVEN CAPABILITY."
    p_win.font.name = 'Space Grotesk'
    p_win.font.size = Pt(18)
    p_win.font.color.rgb = BRAND_ORANGE
    p_win.font.bold = True

    return slide

def add_ryzen():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "AMD Ryzen AI Integration"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Phase 3 Technical Plan"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    code = """CURRENT STATE -> TARGET STATE:
    Cloud LLM Execution  ->  AMD Ryzen AI Distribution

    - All agents on cloud APIs     ->  Planning agents -> CPU (Ryzen)
    - High latency                 ->  Dev agents -> GPU (Radeon)
    - Ongoing cost                 ->  QA agents -> NPU (Ryzen AI)
    - No data control              ->  Local LLM inference
                                     ->  10x efficiency gain

    TECHNICAL REQUIREMENTS:
    1. Add Ryzen AI execution targets to agent definitions
    2. Integrate local LLM (ONNX Runtime / DirectML)
    3. Build hardware-aware agent placement logic
    4. Optimize for NPU tensor operations

    Requires: AMD technical partnership"""

    code_bg = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.8))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = BRAND_GRAY

    tb_code = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.6))
    tf = tb_code.text_frame
    tf.word_wrap = True
    p_code = tf.paragraphs[0]
    p_code.text = code
    p_code.font.name = 'Courier New'
    p_code.font.size = Pt(11)
    p_code.font.color.rgb = BRAND_WHITE

    return slide

def add_roadmap():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "Roadmap (Honest Timeline)"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Phased Delivery with Go/No-Go Gates"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    code = """+------------------------------------------------------------+
|  PHASE 1 (Weeks 1-4): Capture Methodology                  |
|  Document hook patterns | Extract templates | Docs sprint  |
|  [GO] - Entry: User availability confirmed                 |
+------------------------------------------------------------+
|  PHASE 2 (Weeks 5-12): Productization                      |
|  Self-service pipeline | Quality gates | Template library  |
|  [CONDITIONAL GO - extend to 8 weeks if needed]            |
+------------------------------------------------------------+
|  PHASE 3 (Weeks 13-20): Enterprise + Ryzen AI              |
|  Audit trails | AMD optimization | Enterprise packaging    |
|  [GO - AMD partnership dependent]                          |
+------------------------------------------------------------+

    Total: 20 weeks (5 months) to production-ready"""

    code_bg = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.8))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = BRAND_GRAY

    tb_code = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.6))
    tf = tb_code.text_frame
    tf.word_wrap = True
    p_code = tf.paragraphs[0]
    p_code.text = code
    p_code.font.name = 'Courier New'
    p_code.font.size = Pt(11)
    p_code.font.color.rgb = BRAND_WHITE

    return slide

def add_the_ask():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "The Ask"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(36)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "What I Need from AMD"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    # Table
    table = slide.shapes.add_table(5, 3, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.8)).table
    headers = ['Resource', 'Purpose', 'Impact']
    rows = [
        ['Ryzen AI Dev Kit', 'Hardware optimization testing', 'Enable NPU targeting'],
        ['Technical Partnership', 'NPU execution guidance', 'Accelerate Phase 3'],
        ['Go-to-Market Support', 'Enterprise introductions', 'Early adopter pipeline'],
        ['Marketing/PR', 'AMD-optimized positioning', 'Developer adoption']
    ]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_ORANGE
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Archivo Black'
            p.font.size = Pt(11)
            p.font.color.rgb = BRAND_WHITE
            p.font.bold = True
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.name = 'Space Grotesk'
                p.font.size = Pt(11)
                p.font.color.rgb = BRAND_WHITE

    # What AMD Gets
    tb_gets = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(1.2))
    p_gets = tb_gets.text_frame.paragraphs[0]
    p_gets.text = "What AMD Gets:\n- First-mover advantage on hook-based agent systems\n- Differentiated Ryzen AI value proposition\n- Proven capability (Safe Haven running NOW)\n- 41.8K+ star community connection (BMAD)"
    p_gets.font.name = 'Space Grotesk'
    p_gets.font.size = Pt(14)
    p_gets.font.color.rgb = BRAND_WHITE

    return slide

def add_business_case():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "Business Case"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(36)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "ROI Analysis"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    # Stats
    stat_w, stat_h = Inches(3.5), Inches(1.6)
    for i, (num, lbl) in enumerate([
        ('10x', 'Development Speed\nvs. Manual'),
        ('90%+', 'Quality Pass Rate\nProduction-Ready'),
        ('$500K+', 'Annual Savings\nPer Team')
    ]):
        x = Inches(0.5) + i * (stat_w + Inches(0.3))
        card = slide.shapes.add_shape(1, x, Inches(1.2), stat_w, stat_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 100, 50)
        card.line.fill.background()

        tb_num = slide.shapes.add_textbox(x, Inches(1.3), stat_w, Inches(0.8))
        p_num = tb_num.text_frame.paragraphs[0]
        p_num.text = num
        p_num.font.name = 'Archivo Black'
        p_num.font.size = Pt(44)
        p_num.font.color.rgb = BRAND_WHITE
        p_num.alignment = PP_ALIGN.CENTER

        tb_lbl = slide.shapes.add_textbox(x, Inches(2.1), stat_w, Inches(0.6))
        p_lbl = tb_lbl.text_frame.paragraphs[0]
        p_lbl.text = lbl
        p_lbl.font.name = 'Space Grotesk'
        p_lbl.font.size = Pt(12)
        p_lbl.font.color.rgb = BRAND_WHITE
        p_lbl.alignment = PP_ALIGN.CENTER

    # ROI calc
    roi = """ENTERPRISE VALUE:
    Dev time saved: 20 hrs/wk x $100/hr x 50 wks = $100K/dev/year
    Team of 5: $500K/year | Cloud reduction: ~$50K/year | Quality: ~$100K/year
    TOTAL: $650K+/year per enterprise team"""

    tb_roi = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(12.3), Inches(1.2))
    p_roi = tb_roi.text_frame.paragraphs[0]
    p_roi.text = roi
    p_roi.font.name = 'Courier New'
    p_roi.font.size = Pt(12)
    p_roi.font.color.rgb = BRAND_WHITE

    return slide

def add_strategic_approval():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "Strategic Approval"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Adrian-Macias Assessment: APPROVE WITH CONDITIONS"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(16)
    p2.font.color.rgb = BRAND_ORANGE

    # Table
    table = slide.shapes.add_table(6, 3, Inches(0.5), Inches(1.1), Inches(12.3), Inches(2.2)).table
    headers = ['Criterion', 'Rating', 'Notes']
    rows = [
        ['Capability uniqueness', '9/10', 'Genuinely differentiated'],
        ['AMD strategic fit', '8/10', 'Strong Ryzen AI alignment'],
        ['Technical feasibility', '7/10', 'Proven in Safe Haven'],
        ['Business potential', '9/10', '$150B+ TAM, clear ROI'],
        ['OVERALL', 'APPROVE', 'Methodology capture first']
    ]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_ORANGE
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Archivo Black'
            p.font.size = Pt(10)
            p.font.color.rgb = BRAND_WHITE
            p.font.bold = True
    for ri, row in enumerate(rows, 1):
        highlight = (ri == 5)
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            if highlight:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 100, 50)
            for p in cell.text_frame.paragraphs:
                p.font.name = 'Space Grotesk'
                p.font.size = Pt(11)
                p.font.color.rgb = BRAND_DARK if highlight else BRAND_WHITE
                if highlight:
                    p.font.bold = True

    # Conditions
    cond = """Four Mandatory Conditions:
    1. Methodology Capture (Weeks 1-2) - Document context engineering patterns
    2. Reproducibility Validation (Weeks 3-4) - Execute WITHOUT user prompting
    3. AMD Alignment (Weeks 5-6) - Demonstrate Ryzen AI code generation
    4. Enterprise Readiness (Weeks 7-8) - Audit trails, quality gates, monitoring"""

    tb_cond = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.3), Inches(1.5))
    p_cond = tb_cond.text_frame.paragraphs[0]
    p_cond.text = cond
    p_cond.font.name = 'Courier New'
    p_cond.font.size = Pt(11)
    p_cond.font.color.rgb = BRAND_WHITE

    return slide

def add_vision():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "GAIA Vision 2030: Where This Leads"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Outcomes, Possibilities & Predictions"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(18)
    p2.font.color.rgb = BRAND_ORANGE

    # Timeline
    timeline = """PHASE 1 (Year 1): Market Entry
    - 50 enterprise pilots | $8M ARR | 50K developers | AMD partnership

    PHASE 2 (Years 2-3): Category Dominance
    - $65M ARR | 200 Fortune 500 customers | 500K certified developers
    - 15% of AMD Ryzen AI sales driven by GAIA

    PHASE 3 (Years 4-5): Industry Transformation
    - 10x developer productivity standard
    - "Thought to software" in minutes
    - $250M+ ARR | 5M developers | 50% of enterprise software built with GAIA"""

    code_bg = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.3), Inches(3.5))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = BRAND_GRAY

    tb_code = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(3.3))
    tf = tb_code.text_frame
    tf.word_wrap = True
    p_code = tf.paragraphs[0]
    p_code.text = timeline
    p_code.font.name = 'Courier New'
    p_code.font.size = Pt(12)
    p_code.font.color.rgb = BRAND_WHITE

    # Economic impact
    impact = """ECONOMIC IMPACT: $650K/year/team x 1M teams = $650B annual productivity gain
    AMAZON MOMENT: GAIA for software creation = GitHub for version control"""

    tb_impact = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.0))
    p_impact = tb_impact.text_frame.paragraphs[0]
    p_impact.text = impact
    p_impact.font.name = 'Space Grotesk'
    p_impact.font.size = Pt(14)
    p_impact.font.color.rgb = BRAND_WHITE
    p_impact.font.bold = True

    return slide

def add_summary():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "Summary"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(36)
    p.font.color.rgb = BRAND_ORANGE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.5))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "The One-Sentence Pitch"
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(20)
    p2.font.color.rgb = BRAND_ORANGE

    # Pitch
    pitch = '''"I've built Safe Haven - a production hook system with 8 Python hooks, 25+ quality gates, and ChromaDB memory, plus the recursive iterative pipeline with 99.8% test pass rate. I contributed agent templates to BMAD-METHOD before it hit 41.8K stars. GAIA productizes this architecture and optimizes it for AMD Ryzen AI, giving enterprises 'one prompt -> complete feature' capability that's actually battle-tested."'''

    tb_pitch = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.0))
    tf = tb_pitch.text_frame
    tf.word_wrap = True
    p_pitch = tf.paragraphs[0]
    p_pitch.text = pitch
    p_pitch.font.name = 'Space Grotesk'
    p_pitch.font.size = Pt(16)
    p_pitch.font.color.rgb = BRAND_WHITE

    # Stats
    stat_w, stat_h = Inches(3.5), Inches(1.4)
    y_pos = Inches(3.5)
    for i, (num, lbl) in enumerate([
        ('8', 'Production Hooks\nSafe Haven'),
        ('99.8%', 'Test Pass Rate\n1120/1122'),
        ('41.8K+', 'BMAD Stars\nTemplate Contributor')
    ]):
        x = Inches(0.5) + i * (stat_w + Inches(0.3))
        card = slide.shapes.add_shape(1, x, y_pos, stat_w, stat_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 100, 50)
        card.line.fill.background()

        tb_num = slide.shapes.add_textbox(x, y_pos + Inches(0.1), stat_w, Inches(0.7))
        p_num = tb_num.text_frame.paragraphs[0]
        p_num.text = num
        p_num.font.name = 'Archivo Black'
        p_num.font.size = Pt(40)
        p_num.font.color.rgb = BRAND_WHITE
        p_num.alignment = PP_ALIGN.CENTER

        tb_lbl = slide.shapes.add_textbox(x, y_pos + Inches(0.8), stat_w, Inches(0.5))
        p_lbl = tb_lbl.text_frame.paragraphs[0]
        p_lbl.text = lbl
        p_lbl.font.name = 'Space Grotesk'
        p_lbl.font.size = Pt(12)
        p_lbl.font.color.rgb = BRAND_WHITE
        p_lbl.alignment = PP_ALIGN.CENTER

    return slide

def add_contact():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = "Let's Build This Together"
    p.font.name = 'Archivo Black'
    p.font.size = Pt(48)
    p.font.color.rgb = BRAND_ORANGE
    p.alignment = PP_ALIGN.CENTER

    contact = """Anthony Mikinka
AI Engineer & Agent Ecosystem Creator

GitHub: github.com/antmikinka
BMAD-METHOD: 41.8K+ stars
Safe Haven: conversation-compaction-hooks-safe-haven

Strategic Assessment: Adrian-Macias AI Technology Advisor
Decision: APPROVE WITH CONDITIONS"""

    tb_contact = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.3), Inches(3.5))
    tf = tb_contact.text_frame
    tf.word_wrap = True
    for i, line in enumerate(contact.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Space Grotesk'
        p.font.size = Pt(20)
        p.font.color.rgb = BRAND_WHITE
        p.alignment = PP_ALIGN.CENTER

    return slide

# Build presentation
print("Building GAIA presentation...")
add_title_slide()
print("  Slide 1: Title")
add_executive_summary()
print("  Slide 2: Executive Summary")
add_safe_haven()
print("  Slide 3: Safe Haven Proof")
add_recursive_pipeline()
print("  Slide 4: Recursive Pipeline")
add_pipeline_templates()
print("  Slide 5: Pipeline Templates")
add_agent_categories()
print("  Slide 6: Agent Categories")
add_3d_cube()
print("  Slide 7: 3D Cube Matrix")
add_orchestration_modes()
print("  Slide 8: Orchestration Modes")
add_nexus_architecture()
print("  Slide 9: Nexus Architecture")
add_bmad()
print("  Slide 10: BMAD Credibility")
add_metrics()
print("  Slide 11: Metrics")
add_use_cases()
print("  Slide 12: Use Cases")
add_market_opportunity()
print("  Slide 13: Market Opportunity")
add_competitive()
print("  Slide 14: Competitive Landscape")
add_ryzen()
print("  Slide 15: Ryzen AI Integration")
add_roadmap()
print("  Slide 16: Roadmap")
add_the_ask()
print("  Slide 17: The Ask")
add_business_case()
print("  Slide 18: Business Case")
add_strategic_approval()
print("  Slide 19: Strategic Approval")
add_vision()
print("  Slide 20: Vision 2030")
add_summary()
print("  Slide 21: Summary")
add_contact()
print("  Slide 22: Contact")

# Save
output = r"C:\Users\antmi\gaia-proposal\GAIA_Presentation_AMD_v4_FINAL.pptx"
prs.save(output)
print(f"\nDONE! Saved to: {output}")
print(f"Total slides: {len(prs.slides)}")
