#!/usr/bin/env python3
"""Create GAIA Implementation Plan PowerPoint Presentation"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Brand colors
BRAND_ORANGE = RGBColor(255, 87, 34)
BRAND_DARK = RGBColor(26, 26, 26)
BRAND_WHITE = RGBColor(255, 255, 255)
BRAND_GRAY = RGBColor(45, 45, 45)

def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_DARK

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12), Inches(2))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.name = 'Archivo Black'
    p.font.size = Pt(60)
    p.font.color.rgb = BRAND_WHITE

    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(1.5))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.name = 'Space Grotesk'
    p2.font.size = Pt(24)
    p2.font.color.rgb = BRAND_ORANGE

    return slide

def add_content_slide(title, content_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.name = 'Archivo Black'
    p.font.size = Pt(36)
    p.font.color.rgb = BRAND_ORANGE

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = item.get('text', '')
        para.font.name = 'Space Grotesk'
        para.font.size = Pt(item.get('size', 16))
        para.font.color.rgb = BRAND_WHITE
        if item.get('bold', False):
            para.font.bold = True
        if item.get('level', 0) > 0:
            para.level = item.get('level', 0)

    return slide

def add_code_slide(title, code_text, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.name = 'Space Grotesk'
        p2.font.size = Pt(16)
        p2.font.color.rgb = BRAND_ORANGE

    code_bg = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.0))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = BRAND_GRAY
    code_bg.line.fill.background()

    tb_code = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(4.8))
    tf = tb_code.text_frame
    tf.word_wrap = True
    p_code = tf.paragraphs[0]
    p_code.text = code_text
    p_code.font.name = 'Courier New'
    p_code.font.size = Pt(11)
    p_code.font.color.rgb = BRAND_WHITE

    return slide

def add_table_slide(title, headers, rows, highlight_row=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    table_rows = len(rows) + 1
    table_cols = len(headers)
    table = slide.shapes.add_table(table_rows, table_cols, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.8 * table_rows)).table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_ORANGE
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Archivo Black'
            p.font.size = Pt(10)
            p.font.color.rgb = BRAND_WHITE
            p.font.bold = True

    for ri, row in enumerate(rows, 1):
        highlight = highlight_row == ri - 1
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            if highlight:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 100, 50)
            for p in cell.text_frame.paragraphs:
                p.font.name = 'Space Grotesk'
                p.font.size = Pt(11)
                p.font.color.rgb = BRAND_DARK if highlight else BRAND_WHITE
                if highlight:
                    p.font.bold = True

    return slide

def add_diagram_slide(title, diagram_text, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.name = 'Archivo Black'
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_ORANGE

    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.name = 'Space Grotesk'
        p2.font.size = Pt(16)
        p2.font.color.rgb = BRAND_ORANGE

    code_bg = slide.shapes.add_shape(1, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.2))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = BRAND_GRAY

    tb_code = slide.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.0))
    tf = tb_code.text_frame
    tf.word_wrap = True
    p_code = tf.paragraphs[0]
    p_code.text = diagram_text
    p_code.font.name = 'Courier New'
    p_code.font.size = Pt(9)
    p_code.font.color.rgb = BRAND_WHITE

    return slide

# ============================================================================
# SLIDE 1: TITLE
# ============================================================================
add_title_slide(
    "GAIA Implementation Plan",
    "Technical Architecture & Development Roadmap\nAnthony Mikinka | github.com/antmikinka"
)

# ============================================================================
# SLIDE 2: EXECUTIVE SUMMARY
# ============================================================================
add_content_slide("Executive Summary", [
    {'text': "GAIA (Generalized Agent Intelligence Architecture)", 'level': 0, 'size': 20},
    {'text': "", 'level': 0},
    {'text': "MISSION: Productize 'one prompt → complete feature' capability", 'level': 0, 'size': 18},
    {'text': "", 'level': 0},
    {'text': "PROVEN FOUNDATION:", 'level': 0, 'size': 18, 'bold': True},
    {'text': "• Safe Haven: 8 production hooks running in Claude Code", 'level': 1},
    {'text': "• 99.8% test pass rate (1120/1122 tests)", 'level': 1},
    {'text': "• 41.8K+ BMAD-METHOD stars (template contributor)", 'level': 1},
    {'text': "", 'level': 0},
    {'text': "IMPLEMENTATION: 20-week phased delivery", 'level': 0, 'size': 18},
    {'text': "• Phase 1: Core Pipeline Engine (Weeks 1-4)", 'level': 1},
    {'text': "• Phase 2: Productization (Weeks 5-12)", 'level': 1},
    {'text': "• Phase 3: AMD Ryzen AI Integration (Weeks 13-20)", 'level': 1},
])

# ============================================================================
# SLIDE 3: RECURSIVE ITERATIVE PIPELINE
# ============================================================================
pipeline_diagram = """
┌─────────────────────────────────────────────────────────────────────────┐
│              RECURSIVE ITERATIVE PIPELINE (Auto-Pilot)                  │
├─────────────────────────────────────────────────────────────────────────┤
│  USER GOAL → PLANNING → DEVELOPMENT → QUALITY → DECISION                │
│                                     │                                    │
│         ┌───────────────────────────┴────────────────────────────┐      │
│         │                     QUALITY GATE                        │      │
│         │                                                         │      │
│         │   Score >= Threshold?                                   │      │
│         │                                                         │      │
│         │   YES → @software-program-manager → SHIP ✓             │      │
│         │    │                                                    │      │
│         │   NO  → Defects → Loop back to PLANNING                │      │
│         │         (Unlimited iterations until quality met)        │      │
│         └─────────────────────────────────────────────────────────┘      │
└─────────────────────────────────────────────────────────────────────────┘

    KEY INNOVATION: No artificial max iterations. Continues until quality >= threshold."""

add_diagram_slide("Core Architecture: Recursive Iterative Pipeline", pipeline_diagram, "Quality-Gated Loop Architecture")

# ============================================================================
# SLIDE 4: PIPELINE TEMPLATES
# ============================================================================
headers = ['Template', 'Quality Threshold', 'Agent Sequence', 'Use Case']
rows = [
    ['STANDARD', '90/100', 'Planning → Dev → QA → Manager', 'Features, APIs'],
    ['RAPID', '75/100', 'Planning → Dev → QA', 'Prototypes, MVPs'],
    ['ENTERPRISE', '95/100', 'Planning → Dev → QA → Security → Perf → Mgr', 'Production, Security'],
    ['DOCUMENTATION', '85/100', 'Tech Writer → Reviewer → Editor', 'API docs, guides']
]
add_table_slide("Pipeline Templates", headers, rows)

quality_weights = """
QUALITY WEIGHTS (Standard Template):
├── Code Quality:........ 25%  (Syntax, style, complexity, DRY, SOLID)
├── Requirements:........ 25%  (Feature completeness, edge cases)
├── Testing:............. 20%  (Unit tests, integration tests, coverage)
├── Documentation:....... 15%  (Docstrings, README, API docs)
└── Best Practices:...... 15%  (Security, performance, accessibility)"""

add_code_slide("Quality Scoring System", quality_weights)

# ============================================================================
# SLIDE 5: AGENT CATEGORIES
# ============================================================================
agent_categories = """
AGENT CATEGORIES & STATE-BASED ROUTING:

┌──────────────────────────────────────────────────────────────────────────┐
│ PLANNING:                                                                │
│   ├── planning-analysis-strategist  → General technical planning         │
│   ├── solutions-architect           → Complex system architecture        │
│   ├── api-designer                  → REST/GraphQL API design            │
│   └── database-architect            → Database schema & data modeling    │
├──────────────────────────────────────────────────────────────────────────┤
│ DEVELOPMENT:                                                             │
│   ├── senior-developer              → Full-stack generalist              │
│   ├── frontend-specialist           → React, Vue, Angular UI             │
│   ├── backend-specialist            → Server-side APIs                   │
│   ├── devops-engineer               → CI/CD, infrastructure              │
│   └── data-engineer                 → Data pipelines, ETL, ML            │
├──────────────────────────────────────────────────────────────────────────┤
│ REVIEW:                                                                  │
│   ├── quality-reviewer              → General code quality               │
│   ├── security-auditor              → Security vulnerabilities           │
│   ├── performance-analyst           → Performance optimization           │
│   ├── accessibility-reviewer        → WCAG compliance                    │
│   └── test-coverage-analyzer        → Test quality assessment            │
├──────────────────────────────────────────────────────────────────────────┤
│ MANAGEMENT:                                                              │
│   ├── software-program-manager      → Final approval                     │
│   ├── technical-writer              → Documentation                      │
│   └── release-manager               → Deployment coordination            │
└──────────────────────────────────────────────────────────────────────────┘

    AUTO-SELECTION: Based on task triggers ("api" → api-designer, "security" → security-auditor)"""

add_code_slide("Agent Categories", agent_categories, "Multi-Agent Selection System")

# ============================================================================
# SLIDE 6: IMPLEMENTATION PHASES
# ============================================================================
phases_diagram = """
┌────────────────────────────────────────────────────────────────────┐
│  PHASE 1 (Weeks 1-4): Core Pipeline Engine                         │
│  ██████████                                                        │
│  • pipeline.py - Recursive loop engine                             │
│  • quality_scorer.py - 25+ validation categories                   │
│  • orchestrator.py - Agent coordination                            │
│  • state_manager.py - Cross-phase state injection                  │
│  [GO] - Quality threshold: 90/100                                  │
├────────────────────────────────────────────────────────────────────┤
│  PHASE 2 (Weeks 5-12): Productization                              │
│  ████████████                                                      │
│  • Hook System - Safe Haven style (8 hooks)                        │
│  • Agent Registry - State-based routing                            │
│  • Pipeline Templates - STANDARD/RAPID/ENTERPRISE/DOCUMENTATION    │
│  • Test Suite - Unit, Integration, E2E                             │
│  [CONDITIONAL GO - extend to 8 weeks if needed]                    │
├────────────────────────────────────────────────────────────────────┤
│  PHASE 3 (Weeks 13-20): AMD Ryzen AI Integration                   │
│  ██████████                                                        │
│  • Ryzen AI execution targets - CPU/GPU/NPU distribution           │
│  • Local LLM - ONNX Runtime / DirectML                             │
│  • Hardware-aware agent placement                                  │
│  • Performance benchmarks - Target: 10x efficiency                 │
│  [GO - AMD partnership dependent]                                  │
└────────────────────────────────────────────────────────────────────┘

    Total: 20 weeks (5 months) to production-ready"""

add_diagram_slide("Implementation Phases", phases_diagram, "20-Week Phased Delivery with Go/No-Go Gates")

# ============================================================================
# SLIDE 7: PHASE 1 DELIVERABLES
# ============================================================================
headers = ['Component', 'File', 'Function', 'Quality Target']
rows = [
    ['Pipeline Engine', 'pipeline.py', 'Recursive loop with quality gates', '90/100'],
    ['Quality Scorer', 'quality_scorer.py', '25+ validation categories', '90/100'],
    ['Orchestrator', 'orchestrator.py', 'Agent coordination layer', '90/100'],
    ['State Manager', 'state_manager.py', 'Cross-phase state injection', '90/100'],
    ['Config', 'config.py', 'Configuration management', '90/100']
]
add_table_slide("Phase 1: Core Pipeline Engine Deliverables", headers, rows)

# ============================================================================
# SLIDE 8: PHASE 2 DELIVERABLES
# ============================================================================
headers = ['Component', 'Files', 'Function', 'Quality Target']
rows = [
    ['Hook System', '4 Python hooks', 'Pre/Post validation, monitoring', '90/100'],
    ['Agent Registry', 'agent_registry.py', 'State-based routing', '90/100'],
    ['Templates', '4 YAML configs', 'STANDARD/RAPID/ENTERPRISE/DOC', '85/100'],
    ['Test Suite', 'pytest tests', 'Unit/Integration/E2E coverage', '95/100']
]
add_table_slide("Phase 2: Productization Deliverables", headers, rows)

# ============================================================================
# SLIDE 9: PHASE 3 DELIVERABLES
# ============================================================================
headers = ['Component', 'Technology', 'Target', 'Quality Target']
rows = [
    ['Ryzen AI Integration', 'CPU/GPU/NPU', 'Agent distribution', '90/100'],
    ['Local LLM', 'ONNX/DirectML', '<100ms latency', '90/100'],
    ['Hardware Routing', 'Placement logic', 'Automatic agent placement', '90/100'],
    ['Benchmarks', 'Performance report', '10x efficiency gain', '90/100']
]
add_table_slide("Phase 3: AMD Ryzen AI Integration Deliverables", headers, rows)

# ============================================================================
# SLIDE 10: QUALITY ASSURANCE PLAN
# ============================================================================
qa_content = [
    {'text': "MULTI-LAYER TESTING STRATEGY:", 'size': 20, 'bold': True},
    {'text': "", 'size': 16},
    {'text': "┌─────────────────────────────────────────────┐", 'size': 14},
    {'text': "│         E2E Tests (10%) - Playwright       │", 'size': 14},
    {'text': "│    Integration Tests (20%) - pytest        │", 'size': 14},
    {'text': "│       Unit Tests (70%) - pytest            │", 'size': 14},
    {'text': "└─────────────────────────────────────────────┘", 'size': 14},
    {'text': "", 'size': 16},
    {'text': "QUALITY GATES:", 'size': 18, 'bold': True},
    {'text': "• Test Pass Rate: 99%+ target, <95% alert", 'size': 16},
    {'text': "• Pipeline Success Rate: 95%+ target", 'size': 16},
    {'text': "• Average Iterations: <3 target, >5 alert", 'size': 16},
    {'text': "• Quality Score: 90+ target, <85 alert", 'size': 16},
    {'text': "• Response Time: <100ms target, >500ms alert", 'size': 16},
]
add_content_slide("Quality Assurance Plan", qa_content)

# ============================================================================
# SLIDE 11: RESOURCE REQUIREMENTS
# ============================================================================
headers = ['Role', 'Phase 1', 'Phase 2', 'Phase 3', 'Total FTE']
rows = [
    ['Technical Writer', 1, 0, 0, 1],
    ['Senior Engineer', 1, 1, 0, 2],
    ['Backend Engineer', 0, 2, 1, 3],
    ['DevOps Engineer', 0, 1, 0, 1],
    ['QA Engineer', 0, 1, 0, 1],
    ['ML Engineer', 0, 0, 2, 2],
    ['TOTAL', '2', '5', '3', '10']
]
add_table_slide("Resource Requirements", headers, rows, highlight_row=6)

# ============================================================================
# SLIDE 12: RISK ANALYSIS
# ============================================================================
headers = ['Risk', 'Probability', 'Impact', 'Mitigation']
rows = [
    ['Ryzen AI underperforms', 'Low', 'Medium', 'Multi-hardware support roadmap'],
    ['Quality gates fail at scale', 'Low', 'High', 'Continuous validation, pilots'],
    ['Hook system does not generalize', 'Medium', 'High', 'Test across domains early'],
    ['Microsoft/Google copy', 'Medium', 'High', 'First-mover + AMD moat'],
    ['Open source clones', 'High', 'Medium', 'Community building (BMAD model)']
]
add_table_slide("Risk Analysis & Mitigation", headers, rows)

# ============================================================================
# SLIDE 13: SUCCESS METRICS
# ============================================================================
headers = ['Metric', 'Target', 'Measurement', 'Alert']
rows = [
    ['Test Pass Rate', '99%+', 'Automated testing', '<95%'],
    ['Quality Score', '90+', '25+ category validation', '<85'],
    ['Efficiency Gain', '10x', 'Ryzen AI vs. cloud', '<5x'],
    ['Pipeline Success', '95%+', 'Success rate tracking', '<90%'],
    ['Avg Iterations', '<3', 'Loop count tracking', '>5']
]
add_table_slide("Success Metrics & KPIs", headers, rows)

# ============================================================================
# SLIDE 14: TIMELINE & MILESTONES
# ============================================================================
timeline = """
┌────────────────────────────────────────────────────────────────────┐
│  GAIA DEVELOPMENT TIMELINE (20 Weeks)                              │
├────────────────────────────────────────────────────────────────────┤
│                                                                     │
│  WEEK 1-4:    PHASE 1 - Core Pipeline Engine          [M1: GO]     │
│               ████████                                              │
│                                                                     │
│  WEEK 5-12:   PHASE 2 - Productization               [M2: GO]       │
│               ████████████████                                      │
│                                                                     │
│  WEEK 13-20:  PHASE 3 - AMD Ryzen AI Integration     [M3: GO]       │
│               ████████████████                                      │
│                                                                     │
├────────────────────────────────────────────────────────────────────┤
│  MILESTONES:                                                       │
│  M1 (Week 4):  Methodology captured, reproducible                  │
│  M2 (Week 12): 90+ quality score on 10+ agents                     │
│  M3 (Week 20): 10x efficiency, production-ready                    │
└────────────────────────────────────────────────────────────────────┘"""

add_diagram_slide("Timeline & Milestones", timeline, "Key Go/No-Go Gates at Weeks 4, 12, 20")

# ============================================================================
# SLIDE 15: RECURSIVE LOOP IMPLEMENTATION
# ============================================================================
loop_diagram = """
┌─────────────────────────────────────────────────────────────────────────┐
│  IMPLEMENTATION LOOP STRUCTURE (Quality-Gated)                          │
├─────────────────────────────────────────────────────────────────────────┤
│                                                                          │
│  LOOP 1: Core Pipeline Engine                                           │
│  Planning (architect) → Dev (implement) → QA (review) → Decision       │
│  [Quality Gate: 90/100 → Ship or Loop]                                  │
│                                                                          │
│  LOOP 2: Quality Scorer                                                 │
│  Planning (define 25+ categories) → Dev (scorer.py) → QA (validate)    │
│  [Quality Gate: 90/100 → Ship or Loop]                                  │
│                                                                          │
│  LOOP 3: Agent Registry                                                 │
│  Planning (catalog structure) → Dev (registry.py) → QA (routing test)  │
│  [Quality Gate: 90/100 → Ship or Loop]                                  │
│                                                                          │
│  LOOP 4: Hook System                                                    │
│  Planning (Safe Haven patterns) → Dev (4 hooks) → QA (integration)     │
│  [Quality Gate: 90/100 → Ship or Loop]                                  │
│                                                                          │
│  LOOP 5: Templates & Config                                             │
│  Planning (template specs) → Dev (YAML configs) → QA (validation)      │
│  [Quality Gate: 85/100 → Ship or Loop]                                  │
│                                                                          │
│  LOOP 6: Test Suite                                                     │
│  Planning (test strategy) → Dev (pytest tests) → QA (coverage check)   │
│  [Quality Gate: 95/100 → Ship or Loop]                                  │
│                                                                          │
└─────────────────────────────────────────────────────────────────────────┘"""

add_diagram_slide("Recursive Loop Implementation", loop_diagram, "Each Component Goes Through Quality-Gated Loop")

# ============================================================================
# SLIDE 16: AGENT WORKFLOW FOR IMPLEMENTATION
# ============================================================================
workflow = """
┌─────────────────────────────────────────────────────────────────────────┐
│  AGENT WORKFLOW FOR GAIA IMPLEMENTATION                                 │
├─────────────────────────────────────────────────────────────────────────┤
│                                                                          │
│  1. @planning-analysis-strategist (Dr. Sarah Kim)                       │
│     → Technical planning, architecture design, requirements analysis    │
│     → Output: Implementation plan, file structure, component specs      │
│                                                                          │
│     ↓ [Pass to Development]                                              │
│                                                                          │
│  2. @senior-developer                                                   │
│     → Code implementation, component development                        │
│     → Output: Working code, unit tests, documentation                   │
│                                                                          │
│     ↓ [Pass to Quality]                                                  │
│                                                                          │
│  3. @quality-reviewer                                                   │
│     → 25+ category validation, quality scoring                          │
│     → Output: Quality report, defect list, score (0-100)                │
│                                                                          │
│     ↓ [Decision: Score >= 90?]                                           │
│                                                                          │
│     YES → @software-program-manager → SHIP ✓                            │
│     NO  → Loop back to @planning-analysis-strategist (with defects)     │
│                                                                          │
└─────────────────────────────────────────────────────────────────────────┘"""

add_diagram_slide("Agent Workflow", workflow, "Recursive Iterative Loop with Quality Gates")

# ============================================================================
# SLIDE 17: AMD PARTNERSHIP ASK
# ============================================================================
headers = ['Resource', 'Purpose', 'Impact']
rows = [
    ['Ryzen AI Dev Kit', 'Hardware optimization testing', 'Enable NPU targeting'],
    ['Technical Partnership', 'NPU execution guidance', 'Accelerate Phase 3'],
    ['Go-to-Market Support', 'Enterprise introductions', 'Early adopter pipeline'],
    ['Marketing/PR', 'AMD-optimized positioning', 'Developer adoption']
]
add_table_slide("AMD Partnership Requirements", headers, rows)

# ============================================================================
# SLIDE 18: BUSINESS CASE
# ============================================================================
business_content = [
    {'text': "ROI ANALYSIS:", 'size': 20, 'bold': True},
    {'text': "", 'size': 16},
    {'text': "Developer time saved: 20 hrs/week × $100/hr × 50 weeks = $100K/dev/year", 'size': 14},
    {'text': "Team of 5 developers: $500K/year savings", 'size': 14},
    {'text': "Cloud cost reduction (local LLM): ~$50K/year", 'size': 14},
    {'text': "Quality improvement (fewer bugs): ~$100K/year", 'size': 14},
    {'text': "", 'size': 16},
    {'text': "TOTAL: $650K+/year per enterprise team", 'size': 18, 'bold': True},
    {'text': "", 'size': 16},
    {'text': "MARKET OPPORTUNITY:", 'size': 18, 'bold': True},
    {'text': "• Year 1: 50 customers, $8M ARR", 'size': 16},
    {'text': "• Year 2: 200 customers, $65M ARR", 'size': 16},
    {'text': "• Year 3: 500 customers, $250M ARR", 'size': 16},
]
add_content_slide("Business Case & ROI", business_content)

# ============================================================================
# SLIDE 19: NEXT STEPS
# ============================================================================
next_steps = [
    {'text': "IMMEDIATE ACTIONS REQUIRED:", 'size': 20, 'bold': True},
    {'text': "", 'size': 16},
    {'text': "1. APPROVE implementation plan", 'size': 18},
    {'text': "2. Allocate resources (10 FTE across 20 weeks)", 'size': 18},
    {'text': "3. Initiate AMD partnership discussions", 'size': 18},
    {'text': "4. Schedule Phase 1 kickoff (Week 0)", 'size': 18},
    {'text': "", 'size': 16},
    {'text': "WEEK 0 ACTIVITIES:", 'size': 18, 'bold': True},
    {'text': "• Team onboarding", 'size': 16},
    {'text': "• Development environment setup", 'size': 16},
    {'text': "• Repository initialization", 'size': 16},
    {'text': "• Begin Phase 1: Core Pipeline Engine", 'size': 16},
]
add_content_slide("Next Steps", next_steps)

# ============================================================================
# SLIDE 20: SUMMARY & CONTACT
# ============================================================================
summary_text = '''GAIA IMPLEMENTATION SUMMARY

"One Prompt → Complete Feature" Capability

Proven Foundation:
• Safe Haven: 8 production hooks in Claude Code
• 99.8% test pass rate (1120/1122 tests)
• 41.8K+ BMAD-METHOD stars (template contributor)

Implementation: 20 weeks to production-ready
• Phase 1: Core Pipeline Engine (Weeks 1-4)
• Phase 2: Productization (Weeks 5-12)
• Phase 3: AMD Ryzen AI Integration (Weeks 13-20)

Quality Gates: 90/100 threshold with recursive iteration
Resources: 10 FTE across 3 phases
Business Impact: $650K+/year per enterprise team'''

add_content_slide("Summary", [
    {'text': summary_text, 'size': 14}
])

# Contact slide
contact_slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(contact_slide)

tb = contact_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0))
p = tb.text_frame.paragraphs[0]
p.text = "Let's Build This Together"
p.font.name = 'Archivo Black'
p.font.size = Pt(48)
p.font.color.rgb = BRAND_ORANGE
p.alignment = PP_ALIGN.CENTER

contact = """Anthony Mikinka
AI Engineer & Agent Ecosystem Creator

GitHub: github.com/antmikinka
Email: [your email]

Strategic Assessment: Adrian-Macias AI Technology Advisor
Decision: APPROVE WITH CONDITIONS"""

tb_contact = contact_slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.3), Inches(3.5))
tf = tb_contact.text_frame
tf.word_wrap = True
for i, line in enumerate(contact.split('\n')):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = line
    p.font.name = 'Space Grotesk'
    p.font.size = Pt(20)
    p.font.color.rgb = BRAND_WHITE
    p.alignment = PP_ALIGN.CENTER

# Save
output = r"C:\Users\antmi\gaia-proposal\GAIA_IMPLEMENTATION_PLAN.pptx"
prs.save(output)
print(f"DONE! Saved to: {output}")
print(f"Total slides: {len(prs.slides)}")
